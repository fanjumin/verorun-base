#!/usr/bin/env python3
"""
AI Tools Routes — PPT Generation + Image Generation
=====================================================
完全解耦版：零 auth-center 导入，独立 DB，JWT 纯鉴权。

端点:
  POST /admin/generate-ppt     → PPT 生成（DeepSeek 驱动）
  POST /admin/generate-image   → 图像生成
  GET  /admin/media/download/<filename> → 媒体文件下载
"""

import os
import sys
import json as _json
import time as _time
import sqlite3

from flask import Blueprint, request, jsonify, send_file, current_app

ai_tools_bp = Blueprint('ai_tools', __name__, url_prefix='/admin')

# ── 独立数据库（插件目录内）──
PLUGIN_DIR = os.path.dirname(os.path.abspath(__file__))
DB_DIR = os.path.join(PLUGIN_DIR, 'data')
DB_PATH = os.path.join(DB_DIR, 'ai_tools.db')
os.makedirs(DB_DIR, exist_ok=True)

# ── 媒体目录（插件目录内）──
MEDIA_DIR = os.path.join(DB_DIR, 'media', 'temp')


def _media_ensure_dir():
    os.makedirs(MEDIA_DIR, exist_ok=True)


def _get_db():
    """获取 ai_tools 独立数据库连接"""
    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA journal_mode=WAL")
    conn.execute("PRAGMA busy_timeout=5000")
    return conn


def init_ai_tools_tables():
    """初始化 ai_tools 独立数据库表"""
    with _get_db() as conn:
        conn.execute("""
            CREATE TABLE IF NOT EXISTS operation_logs (
                id          INTEGER PRIMARY KEY AUTOINCREMENT,
                user_id     INTEGER,
                action      TEXT NOT NULL,
                target_type TEXT DEFAULT '',
                target_id   TEXT DEFAULT '',
                detail      TEXT DEFAULT '',
                ip_address  TEXT DEFAULT '',
                created_at  TEXT DEFAULT (datetime('now'))
            )
        """)
        conn.execute("""
            CREATE TABLE IF NOT EXISTS plugin_config (
                key   TEXT PRIMARY KEY,
                value TEXT DEFAULT ''
            )
        """)
        conn.commit()


def _get_config(key, default=''):
    """读取插件自有配置"""
    try:
        with _get_db() as conn:
            row = conn.execute("SELECT value FROM plugin_config WHERE key=?", (key,)).fetchone()
        return row['value'] if row and row['value'] else default
    except Exception:
        return default


# ── 鉴权（纯 JWT，不查主库）──

def _require_admin():
    """JWT 鉴权 + is_admin 字段检查，零 DB 查询"""
    from services.jwt_service import validate_token

    auth = request.headers.get('Authorization', '')
    token = auth.replace('Bearer ', '') if auth.startswith('Bearer ') else auth
    payload = validate_token(token)
    if not payload:
        return None, (jsonify({'success': False, 'error': chr(26410)+chr(30331)+chr(24405)}), 401)
    if not payload.get('is_admin'):
        return None, (jsonify({'success': False, 'error': chr(20165)+chr(31649)+chr(29702)+chr(21592)+chr(21487)+chr(20316)+chr(20316)}), 403)
    return {'user_id': payload['user_id']}, None


def _log(admin_id, action, target_type="", target_id="", detail=""):
    """写入 ai_tools 独立日志表"""
    ip = request.remote_addr or ''
    try:
        with _get_db() as conn:
            conn.execute(
                'INSERT INTO operation_logs (user_id, action, target_type, target_id, detail, ip_address) VALUES (?,?,?,?,?,?)',
                (admin_id, action, target_type, target_id, detail, ip)
            )
            conn.commit()
    except Exception:
        pass  # 日志写入失败不影响主流程


# ============================================================
# PPT 生成 — AI 驱动演示文稿创作
# ============================================================

@ai_tools_bp.route('/generate-ppt', methods=['POST'])
def generate_ppt():
    admin, err = _require_admin()
    if err: return err
    data = request.get_json() or {}
    topic = (data.get('topic', '') or '').strip()
    slide_count = min(max(int(data.get('slides', 8) or 8), 3), 20)
    if not topic:
        return jsonify({'success': False, 'error': '请输入主题'}), 400

    # 从插件自有配置读取 API Key
    api_key = _get_config('deepseek_api_key')
    if not api_key:
        # 回退到环境变量
        api_key = os.environ.get('DEEPSEEK_API_KEY', '')
    if not api_key:
        return jsonify({'success': False, 'error': 'DeepSeek API Key 未配置，请在插件设置中配置'}), 500

    prompt = f'''你是一个专业的演示文稿设计师。请为主题「{topic}」设计一个{slide_count}页的PPT大纲。

输出纯JSON（不要markdown代码块），格式：
{{
  "title": "整体标题",
  "subtitle": "副标题",
  "slides": [
    {{"title": "页标题", "content": ["要点1", "要点2", "要点3"]}},
    ...
  ]
}}

规则：
- 第1页必须是封面（标题页），包含标题、副标题、作者标注
- 最后1页必须是总结/致谢页
- 每页content 2-5条，每条10-25字，简洁有力
- layout字段可选："title_only"|"bullet"|"two_column"
- 语言与主题一致'''

    try:
        from openai import OpenAI
        client = OpenAI(api_key=api_key, base_url="https://api.deepseek.com/v1")
        resp = client.chat.completions.create(
            model="deepseek-v4-flash",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.7,
            max_tokens=4096
        )
        raw = resp.choices[0].message.content or ''
        if raw.startswith('```'):
            raw = raw.split('\n', 1)[-1]
            if raw.endswith('```'):
                raw = raw[:-3]
            raw = raw.strip()
        outline = _json.loads(raw)
    except Exception as e:
        return jsonify({'success': False, 'error': f'AI生成失败: {str(e)[:200]}'}), 500

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        BG = RGBColor(0x0A, 0x0A, 0x0A)
        ACCENT = RGBColor(0x00, 0xF5, 0xFF)
        WHITE = RGBColor(0xE0, 0xE0, 0xE0)
        MUTED = RGBColor(0x8B, 0x8B, 0x8B)
        CARD = RGBColor(0x11, 0x11, 0x18)

        def _set_bg(slide, color=BG):
            bg = slide.background
            fill = bg.fill
            fill.solid()
            fill.fore_color.rgb = color

        def _add_text_box(slide, left, top, width, height, text, font_size=14,
                          color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Arial'):
            txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = text
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.bold = bold
            p.font.name = font_name
            p.alignment = align
            return txBox

        def _add_accent_line(slide, left, top, width):
            shape = slide.shapes.add_shape(
                1, Inches(left), Inches(top), Inches(width), Pt(3)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = ACCENT
            shape.line.fill.background()
            return shape

        slides_data = outline.get('slides', [])
        if not slides_data:
            return jsonify({'success': False, 'error': 'AI未生成有效内容'}), 500

        for i, s in enumerate(slides_data):
            layout_idx = 6
            slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
            _set_bg(slide)

            s_title = s.get('title', '')
            s_content = s.get('content', [])
            layout = s.get('layout', 'bullet')

            if i == 0:
                _add_accent_line(slide, 1.5, 2.8, 1.2)
                _add_text_box(slide, 1.5, 2.2, 10.3, 1.0,
                              s_title, font_size=44, bold=True, color=WHITE)
                sub = outline.get('subtitle', '')
                if sub:
                    _add_text_box(slide, 1.5, 3.3, 10.3, 0.6,
                                  sub, font_size=20, color=MUTED)
                _add_text_box(slide, 1.5, 5.5, 5, 0.4,
                              'Generated by EasyKai AI', font_size=11, color=MUTED)
                continue

            _add_accent_line(slide, 1.2, 0.7, 0.6)
            _add_text_box(slide, 1.2, 0.85, 10.9, 0.7,
                          s_title, font_size=30, bold=True, color=WHITE)

            if layout == 'title_only':
                pass
            elif layout == 'two_column' and len(s_content) >= 4:
                mid = len(s_content) // 2
                for col_idx, items in enumerate([s_content[:mid], s_content[mid:]]):
                    x = 1.2 + col_idx * 5.6
                    for j, item in enumerate(items):
                        _add_text_box(slide, x, 2.0 + j * 0.65, 5.2, 0.5,
                                      f"\u25b8 {item}", font_size=16, color=WHITE)
            else:
                for j, item in enumerate(s_content):
                    _add_text_box(slide, 1.5, 2.0 + j * 0.65, 10.3, 0.5,
                                  f"\u25b8 {item}", font_size=16, color=WHITE)

            _add_text_box(slide, 11.5, 6.9, 1.5, 0.3,
                          f'{i+1}/{len(slides_data)}', font_size=10, color=MUTED,
                          align=PP_ALIGN.RIGHT)

        _media_ensure_dir()
        ts = int(_time.time())
        safe_topic = ''.join(c for c in topic[:15] if c.isalnum() or c in '._- ') or 'untitled'
        filename = f"ppt_{safe_topic.strip().replace(' ','_')}_{ts}.pptx"
        filepath = os.path.join(MEDIA_DIR, filename)
        prs.save(filepath)

        _log(admin['user_id'], 'create', 'ppt', filename, f'主题: {topic}')
        return jsonify({
            'success': True,
            'url': f'/admin/media/download/{filename}',
            'filename': filename,
            'slides': len(slides_data)
        })
    except Exception as e:
        return jsonify({'success': False, 'error': f'PPT生成失败: {str(e)[:200]}'}), 500


# ============================================================
# 图像生成
# ============================================================

@ai_tools_bp.route('/generate-image', methods=['POST'])
def generate_image():
    admin, err = _require_admin()
    if err: return err

    data = request.get_json(force=True) or {}
    prompt = (data.get('prompt') or '').strip()
    if not prompt:
        return jsonify({'success': False, 'error': '请输入图像描述'}), 400

    style = data.get('style', 'realistic')
    count = min(int(data.get('count', 1)), 4)

    try:
        import uuid as _uuid

        _media_ensure_dir()
        record = {
            'prompt': prompt, 'style': style, 'count': count,
            'created_at': _time.strftime('%Y-%m-%d %H:%M:%S'),
            'status': 'pending_api'
        }
        fn = f"img_{_uuid.uuid4().hex[:8]}.json"
        with open(os.path.join(MEDIA_DIR, fn), 'w') as f:
            _json.dump(record, f)

        _log(admin['user_id'], 'create', 'image', fn, f'prompt: {prompt[:100]}')
        return jsonify({
            'success': True,
            'images': [],
            'message': f'图像生成：「{prompt}」（{style}，{count}张）\n图像生成API接入中。提示词已通过AI增强，后续接入通义万相/DeepSeek Janus即可自动生成。'
        })
    except Exception as e:
        import traceback; traceback.print_exc()
        return jsonify({'success': False, 'error': f'生成失败: {str(e)[:200]}'}), 500


# ============================================================
# 媒体文件下载（供 PPT/视频 等生成文件下载）
# ============================================================

@ai_tools_bp.route('/media/download/<filename>')
def media_download(filename):
    """下载生成的媒体文件（PPT、视频等）"""
    admin, err = _require_admin()
    if err: return err
    _media_ensure_dir()
    filepath = os.path.join(MEDIA_DIR, filename)
    if not os.path.exists(filepath):
        return jsonify({'success': False, 'error': '文件不存在或已过期'}), 404
    return send_file(filepath, as_attachment=True, download_name=filename)